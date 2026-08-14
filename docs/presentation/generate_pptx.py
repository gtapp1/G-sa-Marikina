"""
Generates G sa Marikina hackathon pitch deck as a .pptx file.
Run: python generate_pptx.py
Output: G-sa-Marikina-Pitch-Deck.pptx (same folder)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---- Brand palette (from DESIGN.md / globals.css) ----
ORANGE = RGBColor(0xE8, 0x59, 0x0C)
ORANGE_DARK = RGBColor(0xC4, 0x32, 0x0A)
RED = RGBColor(0xEB, 0x17, 0x00)
CREAM = RGBColor(0xFF, 0xF8, 0xF3)
BROWN = RGBColor(0x1A, 0x0A, 0x00)
WARM_GRAY = RGBColor(0x6B, 0x3A, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xF0, 0xD0, 0xB8)

FONT = "Sora"
FONT_FALLBACK = "Calibri"  # Sora likely not installed on judges' machines; graceful fallback

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide(bg_color=CREAM):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=BROWN, align=PP_ALIGN.LEFT, font=FONT, italic=False,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_label(slide, left, top, text, color=ORANGE):
    return add_text(slide, left, top, Inches(6), Inches(0.35), text.upper(),
                     size=12, bold=True, color=color)


def add_rounded_card(slide, left, top, width, height, fill=WHITE, line_color=BORDER, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_pill(slide, left, top, width, height, text, fill=ORANGE, font_color=WHITE, size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = font_color
    run.font.name = FONT
    return shape


def add_bullets(slide, left, top, width, height, items, size=13, color=WARM_GRAY, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.2
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    return box


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.9), Inches(1.6), "Hackathon Demo")
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.5),
          "G! sa Marikina", size=64, bold=True, color=BROWN)
add_text(s, Inches(0.9), Inches(3.3), Inches(10), Inches(0.8),
          "Every food spot in your city. One tap away.", size=24, color=WARM_GRAY)
add_text(s, Inches(0.9), Inches(6.6), Inches(10), Inches(0.5),
          "A hyperlocal food discovery platform for Marikina City", size=14, color=WARM_GRAY, italic=True)

# =====================================================================
# SLIDE 2 — THE PAIN
# =====================================================================
s = add_slide(RGBColor(0xFF, 0xF5, 0xF5))
add_label(s, Inches(0.7), Inches(0.5), "The Problem")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(1.1),
          "Small food businesses can't get\nfound online.", size=32, bold=True, color=BROWN)

add_rounded_card(s, Inches(0.7), Inches(2.3), Inches(5.6), Inches(1.6), fill=WHITE)
add_text(s, Inches(1.0), Inches(2.5), Inches(5.0), Inches(0.9),
          '"Selling Ube Crinkles, DM to order!"', size=16, bold=True, color=BROWN)
add_text(s, Inches(1.0), Inches(3.1), Inches(5.0), Inches(0.7),
          "posted 47 times across 12 Facebook groups this week.\nBuried within hours. Every single time.",
          size=12, color=WARM_GRAY)

add_rounded_card(s, Inches(6.7), Inches(2.3), Inches(5.9), Inches(4.0), fill=WHITE)
add_text(s, Inches(7.0), Inches(2.5), Inches(5.2), Inches(0.4),
          "What businesses do today:", size=15, bold=True, color=BROWN)
add_bullets(s, Inches(7.0), Inches(3.0), Inches(5.3), Inches(3.1), [
    "Repost to Facebook groups daily",
    "Posts buried within 2-3 hours",
    "No permanent link to share",
    "Invisible to new customers searching online",
    "No way to collect reviews or build trust",
], size=13)

add_text(s, Inches(0.7), Inches(6.5), Inches(11), Inches(0.6),
          "16 barangays. Hundreds of small food businesses. Zero with a permanent, discoverable URL.",
          size=13, italic=True, color=ORANGE_DARK, bold=True)

# =====================================================================
# SLIDE 3 — WHO SUFFERS
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.7), Inches(0.5), "Two sides of the same problem")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
          "Both sides are stuck.", size=32, bold=True, color=BROWN)

add_rounded_card(s, Inches(0.7), Inches(2.0), Inches(5.7), Inches(4.5), fill=WHITE)
add_text(s, Inches(1.0), Inches(2.25), Inches(5.0), Inches(0.5), "The Business Owner", size=17, bold=True, color=BROWN)
add_bullets(s, Inches(1.0), Inches(2.85), Inches(5.2), Inches(3.4), [
    "Reposts to groups every single day",
    "No website, no search presence",
    "Can't build trust with reviews",
    "Loses customers to competitors with more reach",
    "Spends time marketing instead of cooking",
], size=13)

add_rounded_card(s, Inches(6.9), Inches(2.0), Inches(5.7), Inches(4.5), fill=WHITE)
add_text(s, Inches(7.2), Inches(2.25), Inches(5.0), Inches(0.5), "The Hungry Resident", size=17, bold=True, color=BROWN)
add_bullets(s, Inches(7.2), Inches(2.85), Inches(5.2), Inches(3.4), [
    "Scrolls hundreds of Facebook posts for food",
    "Can't filter by location or category",
    "No reviews to verify quality",
    "Misses great spots two streets away",
    "Relies on word-of-mouth or luck",
], size=13)

# =====================================================================
# SLIDE 4 — THE INSIGHT
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(4.0), Inches(1.4), "The Insight", color=ORANGE)
add_text(s, Inches(1.5), Inches(1.9), Inches(10.3), Inches(1.3),
          "What if every food spot had one URL?", size=34, bold=True, color=BROWN, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.0), Inches(3.3), Inches(9.3), Inches(0.9),
          "Photos. Menu. Map. Reviews. Contact.\nAll in one permanent, shareable page.",
          size=17, color=WARM_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.0), Inches(4.7), Inches(9.3), Inches(1.0),
          "Not another delivery app.\nA food directory for one city.",
          size=18, bold=True, color=BROWN, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 5 — THE SOLUTION
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.7), Inches(0.5), "The Solution")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
          "G! sa Marikina", size=32, bold=True, color=BROWN)
add_text(s, Inches(0.7), Inches(1.55), Inches(10), Inches(0.5),
          "A local food directory built for one city.", size=15, color=WARM_GRAY)

features = [
    ("Browse & Discover", "Categories, barangay filter, map view, curated collections (Top Rated, Hidden Gems, Just Added)"),
    ("Trust Through Reviews", "Real residents leave real ratings — reputation built over time, not one buried post"),
    ("One URL Per Business", "Photos, menu with prices, map pin, phone & Facebook contact, share button"),
    ("Owner Submission", "Submit a spot in 2 minutes. Get listed. Get found. Stop reposting."),
]
positions = [(0.7, 2.3), (6.9, 2.3), (0.7, 4.5), (6.9, 4.5)]
for (title, desc), (l, t) in zip(features, positions):
    add_rounded_card(s, Inches(l), Inches(t), Inches(5.7), Inches(1.9), fill=WHITE)
    add_text(s, Inches(l + 0.3), Inches(t + 0.25), Inches(5.1), Inches(0.5), title, size=16, bold=True, color=BROWN)
    add_text(s, Inches(l + 0.3), Inches(t + 0.85), Inches(5.1), Inches(0.9), desc, size=12, color=WARM_GRAY)

# =====================================================================
# SLIDE 6 — MAP FEATURE
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.7), Inches(0.5), "Standout Feature #1")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
          "Barangay Polygon Map", size=32, bold=True, color=BROWN)
add_text(s, Inches(0.7), Inches(1.6), Inches(6.8), Inches(0.9),
          "Click a barangay. See its real boundaries light up.\nOnly food spots inside appear.",
          size=15, bold=True, color=BROWN)

add_text(s, Inches(0.7), Inches(2.6), Inches(3), Inches(0.4), "How it works:", size=13, bold=True, color=BROWN)
add_bullets(s, Inches(0.7), Inches(3.05), Inches(6.3), Inches(2.4), [
    "Real GeoJSON boundary polygons sourced from OpenStreetMap",
    "Orange fill + outline animates on selection",
    "Camera flies to fit the boundary (MapLibre fitBounds)",
    "Pins outside the boundary disappear — clean, focused",
    "Powered by MapLibre GL + OSM — free, no API key",
], size=12.5)

add_rounded_card(s, Inches(0.7), Inches(5.6), Inches(6.3), Inches(1.3), fill=RGBColor(0xFF, 0xF7, 0xED), line_color=ORANGE)
add_text(s, Inches(1.0), Inches(5.75), Inches(5.7), Inches(1.0),
          'Demo moment: Click "Sta. Elena" → map zooms → Tea Kayo and Aling Beb\'s appear.\nGeographic intelligence in one click.',
          size=12, bold=True, color=ORANGE_DARK)

placeholder = add_rounded_card(s, Inches(7.4), Inches(2.6), Inches(5.2), Inches(4.3), fill=RGBColor(0xEC, 0xEC, 0xEC), line_color=BORDER)
add_text(s, Inches(7.6), Inches(4.6), Inches(4.8), Inches(0.5), "[ Live map demo — screenshot here ]",
          size=12, color=WARM_GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 7 — AI CONCIERGE
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.7), Inches(0.5), "Standout Feature #2")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
          "AI Food Concierge", size=32, bold=True, color=BROWN)
add_text(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.6),
          "Ask in Taglish. Get real answers about real food spots.", size=15, bold=True, color=BROWN)

add_text(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4), "What makes it different:", size=13, bold=True, color=BROWN)
add_bullets(s, Inches(0.7), Inches(2.85), Inches(6.0), Inches(3.0), [
    "Speaks Tagalog, English, or Taglish natively",
    "Only references real listings on the platform",
    "Includes prices, locations, and direct links",
    "Never invents or hallucinates a business",
    "Powered by Gemini 2.5 Flash, streamed responses",
], size=12.5)

add_text(s, Inches(0.7), Inches(5.7), Inches(6.0), Inches(0.8),
          '"This isn\'t ChatGPT with a food skin.\nG! only knows what\'s real."',
          size=13, italic=True, color=WARM_GRAY)

# Chat mockup
chat_l, chat_t, chat_w, chat_h = Inches(7.4), Inches(1.8), Inches(5.2), Inches(5.2)
add_rounded_card(s, chat_l, chat_t, chat_w, chat_h, fill=WHITE, line_color=BORDER)
add_text(s, chat_l + Inches(0.3), chat_t + Inches(0.2), Inches(4.6), Inches(0.4),
          "G! Food Concierge", size=13, bold=True, color=ORANGE)

msg1_bg = add_rounded_card(s, chat_l + Inches(1.3), chat_t + Inches(0.8), Inches(3.6), Inches(0.6), fill=ORANGE)
add_text(s, chat_l + Inches(1.45), chat_t + Inches(0.92), Inches(3.3), Inches(0.4),
          "Saan may masarap na milk tea?", size=11, color=WHITE)

msg2_bg = add_rounded_card(s, chat_l + Inches(0.3), chat_t + Inches(1.55), Inches(4.3), Inches(1.3), fill=RGBColor(0xF5, 0xF0, 0xEB))
add_text(s, chat_l + Inches(0.45), chat_t + Inches(1.65), Inches(4.0), Inches(1.15),
          "Try Tea Kayo Milktea in Sta. Elena! Real tea leaves, not powder. Brown Sugar Fresh Milk is a local favorite at ₱99. /tea-kayo-milktea",
          size=10.5, color=BROWN)

msg3_bg = add_rounded_card(s, chat_l + Inches(0.9), chat_t + Inches(3.0), Inches(4.0), Inches(0.7), fill=ORANGE)
add_text(s, chat_l + Inches(1.05), chat_t + Inches(3.12), Inches(3.7), Inches(0.5),
          "Meron bang mura sa Industrial Valley?", size=11, color=WHITE)

msg4_bg = add_rounded_card(s, chat_l + Inches(0.3), chat_t + Inches(3.85), Inches(4.3), Inches(1.15), fill=RGBColor(0xF5, 0xF0, 0xEB))
add_text(s, chat_l + Inches(0.45), chat_t + Inches(3.95), Inches(4.0), Inches(1.0),
          "Nanay Lita's Karinderya — rice + ulam starts at ₱60. Home-cooked sinigang, adobo, kare-kare. /nanay-litas-karinderya",
          size=10.5, color=BROWN)

# =====================================================================
# SLIDE 8 — TECH STACK
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.7), Inches(0.5), "Under the Hood")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.7),
          "Built lean. Ships fast.", size=30, bold=True, color=BROWN)
add_text(s, Inches(0.7), Inches(1.55), Inches(10), Inches(0.5),
          "Solo developer. All free-tier or open-source tools.", size=14, color=WARM_GRAY)

tech = [
    ("Framework", "Next.js 16"), ("Language", "TypeScript"), ("Styling", "Tailwind v4"),
    ("Maps", "MapLibre GL"), ("Database", "Supabase"), ("Auth", "Clerk"),
    ("ORM", "Drizzle"), ("AI", "Gemini 2.5 Flash"), ("Deploy", "Vercel"),
]
cols, rows = 3, 3
card_w, card_h = Inches(3.7), Inches(1.1)
start_x, start_y = Inches(0.7), Inches(2.3)
gap_x, gap_y = Inches(0.15), Inches(0.2)
for i, (label, name) in enumerate(tech):
    row, col = divmod(i, cols)
    l = start_x + col * (card_w + gap_x)
    t = start_y + row * (card_h + gap_y)
    add_rounded_card(s, l, t, card_w, card_h, fill=WHITE)
    add_text(s, l + Inches(0.2), t + Inches(0.12), card_w - Inches(0.4), Inches(0.3),
              label.upper(), size=10, bold=True, color=WARM_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, l + Inches(0.2), t + Inches(0.5), card_w - Inches(0.4), Inches(0.4),
              name, size=15, bold=True, color=BROWN, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.6), Inches(10), Inches(0.5),
          "Total running cost at demo scale: $0", size=15, bold=True, color=ORANGE)

# =====================================================================
# SLIDE 9 — FLYWHEEL
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.7), Inches(0.5), "Sustainability")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.7),
          "It's a flywheel, not a hack.", size=30, bold=True, color=BROWN)

flow1 = ["Business\nsubmits", "Curator\napproves", "Listed on\nplatform", "Residents\ndiscover"]
flow2 = ["Leave\nreviews", "Business\ngrows", "More businesses\njoin", "Platform gets\nricher"]

def draw_flow(row, items, top):
    n = len(items)
    box_w = Inches(2.6)
    gap = Inches(0.5)
    total_w = n * box_w + (n - 1) * gap
    start_x = (SLIDE_W - total_w) / 2
    for i, label in enumerate(items):
        l = start_x + i * (box_w + gap)
        add_pill(s, l, top, box_w, Inches(0.7), label, fill=WHITE, font_color=BROWN, size=12)
        shp = s.shapes[-1]
        shp.line.color.rgb = ORANGE
        shp.line.width = Pt(1.5)
        if i < n - 1:
            arrow_l = l + box_w
            add_text(s, arrow_l, top, gap, Inches(0.7), "→", size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

draw_flow(1, flow1, Inches(2.1))
draw_flow(2, flow2, Inches(3.1))

add_rounded_card(s, Inches(1.5), Inches(4.5), Inches(4.8), Inches(1.6), fill=WHITE)
add_text(s, Inches(1.8), Inches(4.7), Inches(4.2), Inches(0.4), "Three roles, one loop", size=12, color=WARM_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.8), Inches(5.2), Inches(4.2), Inches(0.7), "Consumer · Business Owner · Curator", size=15, bold=True, color=BROWN, align=PP_ALIGN.CENTER)

add_rounded_card(s, Inches(7.0), Inches(4.5), Inches(4.8), Inches(1.6), fill=WHITE)
add_text(s, Inches(7.3), Inches(4.7), Inches(4.2), Inches(0.6), "More listings = more discovery = more trust", size=12, color=WARM_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.3), Inches(5.35), Inches(4.2), Inches(0.5), "Self-reinforcing growth", size=15, bold=True, color=BROWN, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 10 — LIVE DEMO TRANSITION
# =====================================================================
s = add_slide(BROWN)
add_text(s, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2),
          "Let me show you.", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6),
          "Live demo →", size=18, color=RGBColor(0xE8, 0xC7, 0xB0), align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 11 — IMPACT & VISION
# =====================================================================
s = add_slide(CREAM)
add_label(s, Inches(0.7), Inches(0.5), "Vision")
add_text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(1.2),
          "One city today. Any city tomorrow.", size=30, bold=True, color=BROWN)

stats = [("7", "listings today"), ("30+", "real spots next month"), ("∞", "any city can fork it")]
stat_positions = [1.0, 5.1, 9.2]
for (num, label), x in zip(stats, stat_positions):
    add_text(s, Inches(x), Inches(2.3), Inches(3.2), Inches(1.0), num, size=48, bold=True, color=ORANGE)
    add_text(s, Inches(x), Inches(3.3), Inches(3.2), Inches(0.5), label, size=13, color=WARM_GRAY)
add_text(s, Inches(3.9), Inches(2.5), Inches(1.0), Inches(0.8), "→", size=28, bold=True, color=ORANGE)
add_text(s, Inches(8.0), Inches(2.5), Inches(1.0), Inches(0.8), "→", size=28, bold=True, color=ORANGE)

add_rounded_card(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(2.0), fill=RGBColor(0xFF, 0xF7, 0xED), line_color=ORANGE)
add_text(s, Inches(1.3), Inches(4.8), Inches(10.7), Inches(0.6),
          "The architecture is city-agnostic. Swap the GeoJSON, seed new data, and launch:",
          size=15, bold=True, color=BROWN)
add_text(s, Inches(1.3), Inches(5.5), Inches(10.7), Inches(0.7),
          "G! sa Antipolo · G! sa Pasig · G! sa Cainta", size=18, bold=True, color=ORANGE)

# =====================================================================
# SLIDE 12 — EMOTIONAL CLOSE
# =====================================================================
s = add_slide(CREAM)
add_text(s, Inches(1.2), Inches(1.6), Inches(10.5), Inches(1.6),
          "Nanay Lita has been cooking sinigang\nin Industrial Valley for 15 years.",
          size=28, bold=True, color=BROWN, line_spacing=1.25)
add_text(s, Inches(1.2), Inches(3.6), Inches(9.5), Inches(1.4),
          "She has no website. No Instagram. No Grab listing.\nShe has a karinderya, a loyal neighborhood, and food\nthat tastes like your lola made it.",
          size=16, color=WARM_GRAY, line_spacing=1.4)
add_text(s, Inches(1.2), Inches(5.4), Inches(9.5), Inches(0.8),
          "Now she has a URL.", size=24, bold=True, color=ORANGE)

# =====================================================================
# SLIDE 13 — CLOSING
# =====================================================================
s = add_slide(ORANGE)
add_text(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
          "G! sa Marikina", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.6),
          "Visibility for the invisible food economy.", size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
          "Next.js · Supabase · MapLibre · Gemini AI", size=13, color=RGBColor(0xFF, 0xE8, 0xD6), align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.6),
          "Questions?", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

prs.save("G-sa-Marikina-Pitch-Deck.pptx")
print("Saved G-sa-Marikina-Pitch-Deck.pptx")
